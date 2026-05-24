#!/usr/bin/env python3
"""
build-deck.py — Genera un .pptx draft a partire dall'outline markdown.

USO:
    python scripts/build-deck.py 04-pitch-deck/outline.md output/deck-draft.pptx

DIPENDENZE:
    pip install python-pptx markdown

NOTE:
- Questo è uno scheletro di partenza. Il deck FINALE va rifinito in PowerPoint/Keynote
  con design, immagini, font scelti — questo script genera solo struttura + testo.
- Il parser markdown è minimale: riconosce "## Slide N — Titolo" e usa il primo
  paragrafo come body. Per controllo fine, modificare direttamente il .pptx generato.
- Per design più sofisticato (motivi visivi, palette dedicate), usare la skill
  ufficiale `/mnt/skills/public/pptx/` con pptxgenjs (Node.js).
"""

import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dgm.color import RGBColor
except ImportError as e:
    # Fallback: corregge tipico import errato
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        sys.exit(
            "Errore: pacchetto 'python-pptx' non installato.\n"
            "Installa con: pip install python-pptx"
        )


# Palette sobria: carta, inchiostro, accento ruggine (coerente con identità editoriale)
COLOR_BG = RGBColor(0xFA, 0xF7, 0xF2)        # carta chiara
COLOR_INK = RGBColor(0x1A, 0x1A, 0x1A)       # inchiostro
COLOR_ACCENT = RGBColor(0x8B, 0x3A, 0x2A)    # ruggine

SLIDE_PATTERN = re.compile(r"^##\s+Slide\s+(\d+)\s*[—\-–]\s*(.+?)$", re.MULTILINE)


def parse_outline(md_text: str):
    """Estrae le slide dal markdown nel formato '## Slide N — Titolo'."""
    matches = list(SLIDE_PATTERN.finditer(md_text))
    slides = []

    for i, m in enumerate(matches):
        number = int(m.group(1))
        title = m.group(2).strip().rstrip("*").strip()

        # Body = testo dal match corrente al prossimo
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(md_text)
        body = md_text[start:end].strip()

        # Rimuovi marker markdown di base
        body = re.sub(r"\*\*(.+?)\*\*", r"\1", body)
        body = re.sub(r"\*(.+?)\*", r"\1", body)
        body = re.sub(r"^#+\s+.*$", "", body, flags=re.MULTILINE)
        body = re.sub(r"\n{3,}", "\n\n", body)

        # Prendi solo il primo paragrafo per il body principale (le altre note vanno in speaker notes)
        paragraphs = [p.strip() for p in body.split("\n\n") if p.strip()]
        main_body = paragraphs[0] if paragraphs else ""
        speaker_notes = "\n\n".join(paragraphs[1:]) if len(paragraphs) > 1 else ""

        slides.append({
            "number": number,
            "title": title,
            "body": main_body,
            "notes": speaker_notes,
        })

    return slides


def build_deck(slides, output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # layout vuoto

    for s in slides:
        slide = prs.slides.add_slide(blank_layout)

        # Background colore carta
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

        # Titolo
        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.5), Inches(12), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = f"{s['number']:02d}  {s['title']}"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLOR_INK

        # Body
        body_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.0), Inches(12), Inches(4.5)
        )
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        bp = body_frame.paragraphs[0]
        bp.text = s["body"]
        bp.font.size = Pt(18)
        bp.font.color.rgb = COLOR_INK

        # Footer accento
        footer_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(7.0), Inches(12), Inches(0.3)
        )
        footer_frame = footer_box.text_frame
        fp = footer_frame.paragraphs[0]
        fp.text = "Kalamos AI — Confidenziale — PLAI 2026"
        fp.font.size = Pt(10)
        fp.font.color.rgb = COLOR_ACCENT

        # Speaker notes
        if s["notes"]:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = s["notes"]

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"✓ Deck generato: {output_path}")
    print(f"  Slide totali: {len(slides)}")


def main():
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)

    input_md = Path(sys.argv[1])
    output_pptx = Path(sys.argv[2])

    if not input_md.exists():
        sys.exit(f"Errore: file outline non trovato: {input_md}")

    md_text = input_md.read_text(encoding="utf-8")
    slides = parse_outline(md_text)

    if not slides:
        sys.exit(
            "Errore: nessuna slide trovata nell'outline. "
            "Verifica che le slide siano nel formato '## Slide N — Titolo'."
        )

    build_deck(slides, output_pptx)


if __name__ == "__main__":
    main()
